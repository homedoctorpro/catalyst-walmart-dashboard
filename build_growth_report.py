"""
build_growth_report.py
Builds a 2-sheet Excel + a branded PPTX:
  Sheet/Slide 1: Catalyst subscription growth by year (Shopify DTC, year-end active subs)
  Sheet/Slide 2: Walmart 15-lb Original (CATALYST15ORIG) linear weekly growth — units + U/S/W
"""
import json, glob, os, re, math
from datetime import date, timedelta

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

import openpyxl
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand ─────────────────────────────────────────────────────────────────────
GREEN = "#006633"
INK   = "#1c2b22"
MUTE  = "#6b7d72"
ORANGE = "#E07A2F"
BASE = os.path.dirname(os.path.abspath(__file__))
OUTDIR = os.path.join(BASE, "Growth Report")
os.makedirs(OUTDIR, exist_ok=True)

WALMART_SRC = BASE
FY26_W1_FRIDAY = date(2026, 2, 6)


def safe_float(v):
    try:
        x = float(v)
        return None if math.isnan(x) else x
    except (ValueError, TypeError):
        return None


def week_end_date(code):
    return FY26_W1_FRIDAY + timedelta(days=7 * (int(code[4:]) - 1))


# ── 1. Subscription growth by year (year-end active subs) ──────────────────────
def load_sub_growth():
    series = json.load(open(os.path.join(OUTDIR, "sub_growth_series.json")))
    by_month = {r["month"]: r["active_subscribers"] for r in series}
    years = [2022, 2023, 2024, 2025]                       # matches existing subscriber_growth chart
    vals = [by_month[f"{y}-12"] for y in years]
    return years, vals


# ── 2. Walmart 15O weekly units + U/S/W ────────────────────────────────────────
def find_15o_row(fp):
    """Sheet names vary week to week — scan every sheet for the CATALYST15ORIG row."""
    xl = pd.ExcelFile(fp)
    for sheet in xl.sheet_names:
        try:
            df = pd.read_excel(xl, sheet_name=sheet, header=None, nrows=8)
        except Exception:
            continue
        for i in range(min(6, len(df))):
            if str(df.iloc[i, 0]).strip() == "CATALYST15ORIG" and len(df.columns) > 33:
                return df.iloc[i]
    return None


def load_walmart_15o():
    out = []
    for fp in sorted(glob.glob(os.path.join(WALMART_SRC, "2026?? Weekly Sales Report*.xlsx"))):
        code = os.path.basename(fp)[:6]
        row = find_15o_row(fp)
        if row is None:
            print(f"  [WARN] no 15O row in {code}")
            continue
        units = safe_float(row.iloc[7])
        usw = safe_float(row.iloc[33])
        out.append({
            "week": int(code[4:]),
            "code": code,
            "date": week_end_date(code),
            "units": int(units) if units is not None else None,
            "usw": round(usw, 4) if usw is not None else None,
        })
    out.sort(key=lambda r: r["week"])
    return out


def linfit(x, y):
    """Least-squares slope/intercept; returns (slope, intercept, fitted_array)."""
    x = np.asarray(x, float); y = np.asarray(y, float)
    m, b = np.polyfit(x, y, 1)
    return m, b, m * x + b


# ── Charts (PNG for PPT) ───────────────────────────────────────────────────────
def chart_subscription(years, vals):
    fig, ax = plt.subplots(figsize=(7.6, 5.0), dpi=200)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    x = range(len(years))
    ax.bar(x, vals, width=0.62, color=GREEN, zorder=3)
    for xi, v in zip(x, vals):
        ax.annotate(f"{v:,}", (xi, v), xytext=(0, 7), textcoords="offset points",
                    ha="center", fontsize=12, fontweight="bold", color=INK)
    for i in range(len(years) - 1):
        v0, v1 = vals[i], vals[i + 1]
        pct = (v1 / v0 - 1) * 100
        ax.text(i + 0.5, (v0 + v1) / 2, f"+{pct:,.0f}%", ha="center", va="center",
                fontsize=11, fontweight="bold", color=GREEN,
                bbox=dict(boxstyle="round,pad=0.3", fc="white", ec=GREEN, lw=1.2))
    ax.set_xticks(list(x)); ax.set_xticklabels([str(y) for y in years], fontsize=12, color=INK)
    ax.set_ylim(0, max(vals) * 1.16)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, p: f"{int(v):,}"))
    ax.tick_params(axis="y", labelsize=10, colors=MUTE, length=0)
    ax.tick_params(axis="x", length=0)
    for s in ("top", "right", "left"): ax.spines[s].set_visible(False)
    ax.spines["bottom"].set_color("#d4ddd6")
    ax.grid(axis="y", color="#eef2ef", linewidth=1); ax.set_axisbelow(True)
    fig.tight_layout()
    p = os.path.join(OUTDIR, "rpt_subs.png")
    fig.savefig(p, dpi=200, bbox_inches="tight", facecolor="white"); plt.close(fig)
    return p


def chart_walmart(rows, key, ylabel, fmt, color):
    wk = [r["week"] for r in rows]
    y = [r[key] for r in rows]
    m, b, fit = linfit(wk, y)
    fig, ax = plt.subplots(figsize=(7.6, 5.0), dpi=200)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    ax.plot(wk, y, color=color, lw=2.6, marker="o", ms=5, zorder=3, label="Actual")
    ax.plot(wk, fit, color=INK, lw=1.8, ls="--", zorder=2,
            label=f"Linear trend (+{m:,.2f}/wk)")
    ax.set_xticks(wk)
    ax.set_xticklabels([str(w) for w in wk], fontsize=8.5, color=MUTE)
    ax.set_xlabel("FY26 Walmart week", fontsize=10, color=MUTE)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(fmt))
    ax.tick_params(axis="y", labelsize=10, colors=MUTE, length=0)
    ax.tick_params(axis="x", length=0)
    for s in ("top", "right"): ax.spines[s].set_visible(False)
    ax.spines["left"].set_color("#d4ddd6"); ax.spines["bottom"].set_color("#d4ddd6")
    ax.grid(axis="y", color="#eef2ef", linewidth=1); ax.set_axisbelow(True)
    ax.set_ylabel(ylabel, fontsize=11, color=INK, fontweight="bold")
    ax.legend(frameon=False, fontsize=9.5, loc="upper left")
    fig.tight_layout()
    p = os.path.join(OUTDIR, f"rpt_wm_{key}.png")
    fig.savefig(p, dpi=200, bbox_inches="tight", facecolor="white"); plt.close(fig)
    return p, m, b


# ── Excel ──────────────────────────────────────────────────────────────────────
def hx(c): return c.replace("#", "")

def build_excel(years, subvals, rows, slopes, path):
    wb = openpyxl.Workbook()
    head_fill = PatternFill("solid", fgColor=hx(GREEN))
    head_font = Font(bold=True, color="FFFFFF", size=11)
    title_font = Font(bold=True, color=hx(INK), size=14)
    thin = Side(style="thin", color="D4DDD6")
    border = Border(bottom=thin)

    # Sheet 1 — Subscription Growth
    ws = wb.active; ws.title = "Subscription Growth"
    ws["A1"] = "Catalyst Pet — Subscriber Growth by Year"; ws["A1"].font = title_font
    ws["A2"] = "Active Shopify DTC subscriptions at year-end (excludes ~7,000 Chewy subscribers)"
    ws["A2"].font = Font(italic=True, color=hx(MUTE), size=10)
    hdr = ["Year", "Year-End Active Subscribers", "YoY Growth %"]
    ws.append([]); ws.append(hdr)
    hr = ws.max_row
    for ci, _ in enumerate(hdr, 1):
        c = ws.cell(hr, ci); c.fill = head_fill; c.font = head_font
        c.alignment = Alignment(horizontal="center")
    for i, (y, v) in enumerate(zip(years, subvals)):
        yoy = "" if i == 0 else (v / subvals[i - 1] - 1)
        r = ws.max_row + 1
        ws.cell(r, 1, y)
        ws.cell(r, 2, v).number_format = "#,##0"
        cc = ws.cell(r, 3, yoy)
        if yoy != "":
            cc.number_format = "+0%"; cc.font = Font(bold=True, color=hx(GREEN))
        for ci in range(1, 4): ws.cell(r, ci).border = border
    ws.column_dimensions["A"].width = 10
    ws.column_dimensions["B"].width = 28
    ws.column_dimensions["C"].width = 14

    ch = BarChart(); ch.type = "col"; ch.title = "Subscriber Growth by Year"
    ch.height = 9; ch.width = 17; ch.legend = None
    data = Reference(ws, min_col=2, min_row=hr, max_row=hr + len(years))
    cats = Reference(ws, min_col=1, min_row=hr + 1, max_row=hr + len(years))
    ch.add_data(data, titles_from_data=True); ch.set_categories(cats)
    ch.y_axis.numFmt = "#,##0"; ch.y_axis.majorGridlines = None
    if ch.series:
        from openpyxl.chart.marker import DataPoint
        ch.series[0].graphicalProperties.solidFill = hx(GREEN)
    ws.add_chart(ch, "E4")

    # Sheet 2 — Walmart 15O weekly
    ws2 = wb.create_sheet("Walmart 15O Weekly")
    ws2["A1"] = "Walmart — 15 lb Original (CATALYST15ORIG) Weekly Growth"; ws2["A1"].font = title_font
    ws2["A2"] = ("In-store units & U/S/W (units per store per week), FY26 weekly, "
                 "with least-squares linear trend")
    ws2["A2"].font = Font(italic=True, color=hx(MUTE), size=10)
    hdr2 = ["FY Week", "Week-Ending", "Units", "U/S/W", "Units (Linear Trend)", "U/S/W (Linear Trend)"]
    ws2.append([]); ws2.append(hdr2)
    hr2 = ws2.max_row
    for ci, _ in enumerate(hdr2, 1):
        c = ws2.cell(hr2, ci); c.fill = head_fill; c.font = head_font
        c.alignment = Alignment(horizontal="center")
    (mu, bu), (mw, bw) = slopes["units"], slopes["usw"]
    for r in rows:
        rr = ws2.max_row + 1
        ws2.cell(rr, 1, r["week"])
        ws2.cell(rr, 2, r["date"]).number_format = "mm/dd/yy"
        ws2.cell(rr, 3, r["units"]).number_format = "#,##0"
        ws2.cell(rr, 4, r["usw"]).number_format = "0.000"
        ws2.cell(rr, 5, round(mu * r["week"] + bu, 1)).number_format = "#,##0"
        ws2.cell(rr, 6, round(mw * r["week"] + bw, 4)).number_format = "0.000"
        for ci in range(1, 7): ws2.cell(rr, ci).border = border
    last = ws2.max_row
    # slope summary
    sr = last + 2
    ws2.cell(sr, 1, "Linear weekly growth:").font = Font(bold=True, color=hx(INK))
    ws2.cell(sr, 3, f"+{mu:,.0f} units/wk").font = Font(bold=True, color=hx(GREEN))
    ws2.cell(sr, 4, f"+{mw:.3f} U/S/W per wk").font = Font(bold=True, color=hx(GREEN))
    for col, w in zip("ABCDEF", (9, 13, 12, 10, 20, 20)):
        ws2.column_dimensions[col].width = w

    lc = LineChart(); lc.title = "15O Weekly Units (actual vs linear trend)"
    lc.height = 8.5; lc.width = 17; lc.style = 2
    udata = Reference(ws2, min_col=3, min_row=hr2, max_row=last)
    tdata = Reference(ws2, min_col=5, min_row=hr2, max_row=last)
    cats2 = Reference(ws2, min_col=1, min_row=hr2 + 1, max_row=last)
    lc.add_data(udata, titles_from_data=True); lc.add_data(tdata, titles_from_data=True)
    lc.set_categories(cats2); lc.y_axis.numFmt = "#,##0"
    if len(lc.series) >= 2:
        lc.series[0].graphicalProperties.line.solidFill = hx(ORANGE)
        lc.series[0].graphicalProperties.line.width = 28000
        lc.series[1].graphicalProperties.line.solidFill = hx(INK)
        lc.series[1].graphicalProperties.line.dashStyle = "dash"
    ws2.add_chart(lc, "H4")

    lc2 = LineChart(); lc2.title = "15O Weekly U/S/W (actual vs linear trend)"
    lc2.height = 8.5; lc2.width = 17; lc2.style = 2
    wdata = Reference(ws2, min_col=4, min_row=hr2, max_row=last)
    wtdata = Reference(ws2, min_col=6, min_row=hr2, max_row=last)
    lc2.add_data(wdata, titles_from_data=True); lc2.add_data(wtdata, titles_from_data=True)
    lc2.set_categories(cats2); lc2.y_axis.numFmt = "0.000"
    if len(lc2.series) >= 2:
        lc2.series[0].graphicalProperties.line.solidFill = hx(GREEN)
        lc2.series[0].graphicalProperties.line.width = 28000
        lc2.series[1].graphicalProperties.line.solidFill = hx(INK)
        lc2.series[1].graphicalProperties.line.dashStyle = "dash"
    ws2.add_chart(lc2, "H22")

    wb.save(path)


# ── PPT ──────────────────────────────────────────────────────────────────────
def _bg(slide, prs, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(hx(color))

def build_ppt(sub_png, wm_units_png, wm_usw_png, years, subvals, slopes, rows, path):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height

    def textbox(slide, l, t, w, h, text, size, color, bold=True, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text; f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = RGBColor.from_string(hx(color)); f.name = "Futura"
        return tb

    # Slide 1 — title
    s = prs.slides.add_slide(blank); _bg(s, prs, GREEN)
    textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4),
            "Catalyst Pet — Growth Report", 44, "#FFFFFF")
    textbox(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.8),
            "Subscriber growth by year  •  Walmart 15 lb Original weekly trajectory",
            20, "#CFE6D7", bold=False)
    textbox(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
            f"Generated {date.today():%B %d, %Y}", 12, "#9FC7AE", bold=False)

    # Slide 2 — subscription growth
    s = prs.slides.add_slide(blank); _bg(s, prs, "#FFFFFF")
    textbox(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
            "Subscriber Growth by Year", 30, INK)
    textbox(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
            "Active Shopify DTC subscriptions at year-end (excludes ~7,000 Chewy subscribers)",
            14, MUTE, bold=False, italic=True)
    cagr = (subvals[-1] / subvals[0]) ** (1 / (len(subvals) - 1)) - 1
    s.shapes.add_picture(sub_png, Inches(0.6), Inches(1.8), height=Inches(5.0))
    # callout
    textbox(s, Inches(8.2), Inches(2.4), Inches(4.6), Inches(0.6),
            f"{subvals[0]:,} → {subvals[-1]:,}", 30, GREEN)
    textbox(s, Inches(8.2), Inches(3.2), Inches(4.6), Inches(1.2),
            f"{years[0]}–{years[-1]} year-end\n~{cagr*100:,.0f}% CAGR", 18, INK, bold=False)

    # Slide 3 — Walmart 15O weekly
    s = prs.slides.add_slide(blank); _bg(s, prs, "#FFFFFF")
    textbox(s, Inches(0.6), Inches(0.4), Inches(12.5), Inches(0.8),
            "Walmart — 15 lb Original: Linear Weekly Growth", 28, INK)
    textbox(s, Inches(0.6), Inches(1.15), Inches(12.5), Inches(0.5),
            "In-store units and U/S/W (units per store per week), FY26 weekly with linear trend",
            14, MUTE, bold=False, italic=True)
    s.shapes.add_picture(wm_units_png, Inches(0.35), Inches(1.85), height=Inches(4.5))
    s.shapes.add_picture(wm_usw_png, Inches(6.85), Inches(1.85), height=Inches(4.5))
    mu = slopes["units"][0]; mw = slopes["usw"][0]
    textbox(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.6),
            f"Linear weekly growth:  +{mu:,.0f} units/week   •   +{mw:.3f} U/S/W per week"
            f"   (Wk {rows[0]['week']}: {rows[0]['units']:,} u / {rows[0]['usw']:.3f}  →  "
            f"Wk {rows[-1]['week']}: {rows[-1]['units']:,} u / {rows[-1]['usw']:.3f})",
            14, GREEN)

    prs.save(path)


def sync_inbox():
    """Pull any weekly workbooks from the reports inbox that aren't on disk yet,
    so the report never silently builds against a stale folder."""
    try:
        from _download_missing_weeks import download_missing_weeks
        print("Syncing weekly workbooks from inbox …")
        download_missing_weeks()
    except Exception as e:
        print(f"  [WARN] inbox sync failed ({e}) — building with workbooks on disk")


def main():
    sync_inbox()
    print("Loading subscription growth …")
    years, subvals = load_sub_growth()
    print("  ", dict(zip(years, subvals)))
    print("Loading Walmart 15O weekly …")
    rows = load_walmart_15o()
    print(f"   {len(rows)} weeks, wk{rows[0]['week']}->wk{rows[-1]['week']}")

    mu, bu, _ = linfit([r["week"] for r in rows], [r["units"] for r in rows])
    mw, bw, _ = linfit([r["week"] for r in rows], [r["usw"] for r in rows])
    slopes = {"units": (mu, bu), "usw": (mw, bw)}
    print(f"   units slope +{mu:,.1f}/wk ; U/S/W slope +{mw:.4f}/wk")

    sub_png = chart_subscription(years, subvals)
    wm_u_png, _, _ = chart_walmart(rows, "units", "In-store units",
                                   lambda v, p: f"{int(v):,}", ORANGE)
    wm_w_png, _, _ = chart_walmart(rows, "usw", "U/S/W",
                                   lambda v, p: f"{v:.2f}", GREEN)

    xlsx = os.path.join(OUTDIR, "Catalyst_Growth_Report.xlsx")
    pptx = os.path.join(OUTDIR, "Catalyst_Growth_Report.pptx")
    build_excel(years, subvals, rows, slopes, xlsx)
    build_ppt(sub_png, wm_u_png, wm_w_png, years, subvals, slopes, rows, pptx)
    print("Saved:")
    print("  ", xlsx)
    print("  ", pptx)


if __name__ == "__main__":
    main()
